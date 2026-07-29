"""
Documents Clients — Extracteur de texte multi-format.

Supporte : PDF, DOCX, TXT, CSV, PPTX
Applique un nettoyage de base (espaces multiples, lignes vides, caractères de contrôle).

Dépendances optionnelles — si une lib est absente, on tente quand même les autres :
  - pdfminer.six  → PDF
  - python-docx   → DOCX
  - python-pptx   → PPTX
  (CSV et TXT n'ont pas de dépendances extras)
"""

import csv
import io
import logging
import re
import unicodedata

logger = logging.getLogger(__name__)

# ── Nettoyage ─────────────────────────────────────────────────────────────────

_CTRL_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")   # caractères de contrôle
_MULTI_SPACE_RE = re.compile(r"[ \t]+")
_MULTI_NL_RE = re.compile(r"\n{3,}")


def clean_text(raw: str) -> str:
    """
    Nettoyage standard du texte extrait :
    1. Normalisation Unicode (NFC)
    2. Suppression des caractères de contrôle
    3. Collapsing des espaces/tabulations multiples
    4. Collapsing des lignes vides multiples (max 2 consécutives)
    5. Strip final
    """
    text = unicodedata.normalize("NFC", raw)
    text = _CTRL_RE.sub(" ", text)
    text = _MULTI_SPACE_RE.sub(" ", text)
    text = _MULTI_NL_RE.sub("\n\n", text)
    return text.strip()


# ── Extracteurs par type ───────────────────────────────────────────────────────

def _extract_txt(content: bytes) -> str:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _extract_csv(content: bytes) -> str:
    text = _extract_txt(content)
    reader = csv.reader(io.StringIO(text))
    lines = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
    return "\n".join(lines)


def _extract_pdf(content: bytes) -> str:
    try:
        from pdfminer.high_level import extract_text_to_fp
        from pdfminer.layout import LAParams
        output = io.StringIO()
        extract_text_to_fp(io.BytesIO(content), output, laparams=LAParams())
        return output.getvalue()
    except ImportError:
        logger.warning("[Extractor] pdfminer.six non installé — PDF ignoré")
        return ""
    except Exception as e:
        logger.error(f"[Extractor] Erreur extraction PDF : {e}")
        return ""


def _extract_docx(content: bytes) -> str:
    try:
        import docx
        doc = docx.Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Inclure les tableaux
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs)
    except ImportError:
        logger.warning("[Extractor] python-docx non installé — DOCX ignoré")
        return ""
    except Exception as e:
        logger.error(f"[Extractor] Erreur extraction DOCX : {e}")
        return ""


def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"[Diapositive {i}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            slides_text.append("\n".join(parts))
        return "\n\n".join(slides_text)
    except ImportError:
        logger.warning("[Extractor] python-pptx non installé — PPTX ignoré")
        return ""
    except Exception as e:
        logger.error(f"[Extractor] Erreur extraction PPTX : {e}")
        return ""


# ── Point d'entrée principal ──────────────────────────────────────────────────

EXTRACTOR_MAP = {
    ".txt":  _extract_txt,
    ".csv":  _extract_csv,
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".doc":  _extract_docx,
    ".pptx": _extract_pptx,
    ".ppt":  _extract_pptx,
}


def extract_and_clean(filename: str, content: bytes) -> str:
    """
    Extrait le texte brut d'un fichier selon son extension, puis le nettoie.
    Retourne une chaîne vide si le format n'est pas supporté.
    """
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    extractor = EXTRACTOR_MAP.get(ext)
    if extractor is None:
        logger.warning(f"[Extractor] Format non supporté : {ext!r}")
        return ""
    raw = extractor(content)
    cleaned = clean_text(raw)
    logger.info(f"[Extractor] {filename} ({ext}) → {len(cleaned)} caractères après nettoyage")
    return cleaned
